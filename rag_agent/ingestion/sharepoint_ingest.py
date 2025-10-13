# # src/ingestion/sharepoint_ingest.py
# from typing import List, Dict, Any
# import asyncio
# from office365.runtime.auth.authentication_context import AuthenticationContext
# from office365.sharepoint.client_context import ClientContext
# from office365.sharepoint.files.file import File
# import PyPDF2
# from docx import Document
# from pptx import Presentation
# import io
# import os
# from langchain_text_splitters.character import RecursiveCharacterTextSplitter
# from langchain_community.vectorstores.azuresearch import AzureSearch
# from rag_agent.LLM_Garden import embeddings, LLM_GPT_5_MINI
# from dotenv import load_dotenv
# load_dotenv()
# AZURE_SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT")
# AZURE_SEARCH_KEY = os.getenv("AZURE_SEARCH_KEY")


# class SharePointIngestion:
#     def __init__(self, index_name: str):
#         self.site_url = os.getenv("SHAREPOINT_URL")
#         self.client_id = os.getenv("SHAREPOINT_CLIENT_ID")
#         self.client_secret = os.getenv("SHAREPOINT_CLIENT_SECRET_ID_VALUE")
#         self.azure_ai_search_client =  AzureSearch(
#                                         azure_search_endpoint=AZURE_SEARCH_ENDPOINT,
#                                         azure_search_key=AZURE_SEARCH_KEY,
#                                         index_name=self.index_name,
#                                         embedding_function=embeddings.embed_query,
#                                     )
        
#         self.embedding_model = embeddings
#         self.chunk_size = 1000
#         self.chunk_overlap = 200
#         self.text_splitter = RecursiveCharacterTextSplitter(
#                                 # Set a really small chunk size, just to show.
#                                 chunk_size=self.chunk_size,
#                                 chunk_overlap=self.chunk_overlap,
#                                 length_function=len,
#                                 is_separator_regex=False,
#                             )
#         self.index_name = index_name
#     async def ingest_documents(self, document_library_path: str) -> Dict[str, Any]:
#         """Main ingestion pipeline for SharePoint documents"""
#         # Authenticate
#         ctx = self._authenticate()
        
#         # Get documents
#         documents = self._get_documents(ctx, document_library_path)
        
#         # Process documents in parallel
#         tasks = []
#         for doc in documents:
#             tasks.append(self._process_document(doc))
        
#         results = await asyncio.gather(*tasks)
        
#         # Index to OpenSearch
#         indexed_count = await self._index_documents(results)
        
#         return {
#             'total_documents': len(documents),
#             'indexed_count': indexed_count,
#             'errors': self._collect_errors(results)
#         }
    
#     async def _process_document(self, doc_info: Dict[str, Any]) -> Dict[str, Any]:
#         """Process individual document"""
#         file_content = self._download_file(doc_info['url'])
        
#         # Extract text based on file type
#         if doc_info['extension'] == '.pdf':
#             text = self._extract_pdf_text(file_content)
#         elif doc_info['extension'] in ['.docx', '.doc']:
#             text = self._extract_word_text(file_content)
#         elif doc_info['extension'] in ['.pptx', '.ppt']:
#             text = self._extract_powerpoint_text(file_content)
#         else:
#             text = file_content.decode('utf-8', errors='ignore')
        
#         # Chunk text
#         chunks = self._chunk_text(text)
        
#         # Generate embeddings for each chunk
#         processed_chunks = []
#         for i, chunk in enumerate(chunks):
#             embedding = await self._generate_embedding(chunk)
            
#             processed_chunks.append({
#                 'doc_id': f"{doc_info['id']}_{i}",
#                 'source': 'SharePoint',
#                 'file_type': doc_info['extension'],
#                 'page_number': i,
#                 'title': doc_info['name'],
#                 'embedding': embedding,
#                 'content': chunk,
#                 'summary': await self._generate_summary(chunk),
#                 'metadata': {
#                     'author': doc_info.get('author'),
#                     'modified': doc_info.get('modified'),
#                     'url': doc_info['url']
#                 }
#             })
        
#         return processed_chunks
    
#     def _chunk_text(self, text: str) -> List[str]:
#         """Split text into overlapping chunks"""
#         return self.text_splitter.split_text(text)
    
#     async def _index_documents(self, document_chunks: List[List[Dict]]) -> int:
#         """Index processed documents to OpenSearch"""
#         actions = []
        
#         for doc_chunks in document_chunks:
#             for chunk in doc_chunks:
#                 actions.append({
#                     "index": {
#                         "_index": self.index_name,
#                         "_id": chunk['doc_id']
#                     }
#                 })
#                 actions.append(chunk)
        
#         # Bulk index
#         response = self.azure_ai_search_client.add_texts(documents=actions)
        
#         return len([item for item in response['items'] if item['index']['result'] == 'created'])
    
    # src/ingestion/sharepoint_ingest.py
from typing import List, Dict, Any
import asyncio
from office365.runtime.auth.authentication_context import AuthenticationContext
from office365.sharepoint.client_context import ClientContext
from office365.sharepoint.files.file import File
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import io
import os
from langchain_text_splitters.character import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.azuresearch import AzureSearch
from rag_agent.LLM_Garden import embeddings,LLM_GPT_5_MINI
from dotenv import load_dotenv

load_dotenv()
AZURE_SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT")
AZURE_SEARCH_KEY = os.getenv("AZURE_SEARCH_KEY")


class SharePointIngestion:
    def __init__(self, index_name: str):
        self.site_url = os.getenv("SHAREPOINT_URL")
        self.client_id = os.getenv("SHAREPOINT_CLIENT_ID")
        self.client_secret = os.getenv("SHAREPOINT_CLIENT_SECRET_ID_VALUE")
        self.index_name = index_name
        
        self.azure_ai_search_client = AzureSearch(
            azure_search_endpoint=AZURE_SEARCH_ENDPOINT,
            azure_search_key=AZURE_SEARCH_KEY,
            index_name=self.index_name,
            embedding_function=embeddings.embed_query,
        )
        
        self.embedding_model = embeddings
        self.chunk_size = 1000
        self.chunk_overlap = 200
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            is_separator_regex=False,
        )
        
    def _authenticate(self) -> ClientContext:
        """Authenticate with SharePoint using app credentials"""
        context_auth = AuthenticationContext(url=self.site_url)
        context_auth.acquire_token_for_app(
            client_id=self.client_id,  
            client_secret=self.client_secret
        )
        ctx = ClientContext(self.site_url, context_auth)
        return ctx
    
    def _get_documents(self, ctx: ClientContext, document_library_path: str) -> List[Dict[str, Any]]:
        """Get list of documents from SharePoint library"""
        documents = []
        
        # Get the folder
        folder = ctx.web.get_folder_by_server_relative_url(document_library_path)
        folder.expand(["Files"]).get().execute_query()
        
        # Process each file
        for file in folder.files:
            file_extension = os.path.splitext(file.name)[1].lower()
            documents.append({
                'id': file.unique_id,
                'name': file.name,
                'extension': file_extension,
                'url': file.serverRelativeUrl,
                'author': file.author.title if hasattr(file, 'author') else 'Unknown',
                'modified': file.time_last_modified.isoformat() if hasattr(file, 'time_last_modified') else None
            })
        
        return documents
    
    def _download_file(self, file_url: str) -> bytes:
        """Download file content from SharePoint"""
        ctx = self._authenticate()
        response = File.open_binary(ctx, file_url)
        return response.content
    
    def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF using PyMuPDF (fitz)"""
        text = ""
        try:
            # Open PDF from bytes
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            
            # Extract text from each page
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text += page.get_text()
            
            pdf_document.close()
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            text = ""
        
        return text
    
    def _extract_word_text(self, file_content: bytes) -> str:
        """Extract text from Word documents"""
        text = ""
        try:
            doc = Document(io.BytesIO(file_content))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                text += "\n"
        except Exception as e:
            print(f"Error extracting Word text: {e}")
            text = ""
        
        return text
    
    def _extract_powerpoint_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        text = ""
        try:
            prs = Presentation(io.BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting PowerPoint text: {e}")
            text = ""
        
        return text
    
    async def _generate_embedding(self, text: str) -> List[float]:
        """Generate embedding for text chunk"""
        try:
            # Use the embedding model's embed_query method
            embedding = self.embedding_model.embed_query(text)
            return embedding
        except Exception as e:
            print(f"Error generating embedding: {e}")
            return []
    
    async def _generate_summary(self, text: str) -> str:
        """Generate summary for text chunk using LLM"""
        try:
            # Limit text length for summary generation
            text_preview = text[:500] if len(text) > 500 else text
            
            prompt = f"Summarize the following text in 1-2 sentences:\n\n{text_preview}"
            
            # Use the chat model from LLM_Garden
            summary = await LLM_GPT_5_MINI.ainvoke(prompt)
            
            # Extract content from response
            if hasattr(summary, 'content'):
                return summary.content
            else:
                return str(summary)
        except Exception as e:
            print(f"Error generating summary: {e}")
            return text[:200]  # Fallback to first 200 chars
    
    def _collect_errors(self, results: List[List[Dict]]) -> List[str]:
        """Collect any errors from processing results"""
        errors = []
        for result in results:
            if isinstance(result, dict) and 'error' in result:
                errors.append(result['error'])
        return errors
    
    async def ingest_documents(self, document_library_path: str) -> Dict[str, Any]:
        """Main ingestion pipeline for SharePoint documents"""
        # Authenticate
        ctx = self._authenticate()
        
        # Get documents
        documents = self._get_documents(ctx, document_library_path)
        
        # Process documents in parallel
        tasks = []
        for doc in documents:
            tasks.append(self._process_document(doc))
        
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Filter out exceptions
        valid_results = [r for r in results if not isinstance(r, Exception)]
        
        # Index to Azure AI Search
        indexed_count = await self._index_documents(valid_results)
        
        return {
            'total_documents': len(documents),
            'indexed_count': indexed_count,
            'errors': self._collect_errors(results)
        }
    
    async def _process_document(self, doc_info: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Process individual document"""
        try:
            file_content = self._download_file(doc_info['url'])
            
            # Extract text based on file type
            if doc_info['extension'] == '.pdf':
                text = self._extract_pdf_text(file_content)
            elif doc_info['extension'] in ['.docx', '.doc']:
                text = self._extract_word_text(file_content)
            elif doc_info['extension'] in ['.pptx', '.ppt']:
                text = self._extract_powerpoint_text(file_content)
            else:
                text = file_content.decode('utf-8', errors='ignore')
            
            # Chunk text
            chunks = self._chunk_text(text)
            
            # Generate embeddings and summaries for each chunk
            processed_chunks = []
            for i, chunk in enumerate(chunks):
                summary = await self._generate_summary(chunk)
                
                processed_chunks.append({
                    'doc_id': f"{doc_info['id']}_{i}",
                    'source': 'SharePoint',
                    'file_type': doc_info['extension'],
                    'page_number': i,
                    'title': doc_info['name'],
                    'content': chunk,
                    'summary': summary,
                    'metadata': {
                        'author': doc_info.get('author'),
                        'modified': doc_info.get('modified'),
                        'url': doc_info['url'],
                        'chunk_index': i,
                        'total_chunks': len(chunks)
                    }
                })
            
            return processed_chunks
        except Exception as e:
            print(f"Error processing document {doc_info.get('name', 'unknown')}: {e}")
            return []
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks"""
        return self.text_splitter.split_text(text)
    
    async def _index_documents(self, document_chunks: List[List[Dict]]) -> int:
        """Index processed documents to Azure AI Search with proper bulk upload"""
        # Flatten the list of chunks
        all_chunks = []
        for doc_chunks in document_chunks:
            if isinstance(doc_chunks, list):
                all_chunks.extend(doc_chunks)
        
        if not all_chunks:
            return 0
        
        # Prepare texts and metadatas for bulk upload
        texts = []
        metadatas = []
        keys = []
        
        for chunk in all_chunks:
            texts.append(chunk['content'])
            
            # Prepare metadata dictionary
            metadata = {
                'source': chunk['source'],
                'file_type': chunk['file_type'],
                'page_number': chunk['page_number'],
                'title': chunk['title'],
                'summary': chunk['summary'],
                'author': chunk['metadata'].get('author', 'Unknown'),
                'modified': chunk['metadata'].get('modified', ''),
                'url': chunk['metadata'].get('url', ''),
                'chunk_index': chunk['metadata'].get('chunk_index', 0),
                'total_chunks': chunk['metadata'].get('total_chunks', 0)
            }
            metadatas.append(metadata)
            keys.append(chunk['doc_id'])
        
        try:
            # Use add_texts method with proper parameters
            # The AzureSearch vectorstore handles embedding generation internally
            ids = self.azure_ai_search_client.add_texts(
                texts=texts,
                metadatas=metadatas,
                keys=keys
            )
            
            return len(ids)
        except Exception as e:
            print(f"Error indexing documents to Azure AI Search: {e}")
            return 0

import os
import asyncio
import io
from typing import List, Dict, Any
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.azuresearch import AzureSearch

class LocalFileIngestion:
    def __init__(self, index_name: str):
        self.index_name = index_name
        
        # Azure AI Search configuration
        self.azure_ai_search_client = AzureSearch(
            azure_search_endpoint=AZURE_SEARCH_ENDPOINT,
            azure_search_key=AZURE_SEARCH_KEY,
            index_name=self.index_name,
            embedding_function=embeddings,
        )
        
        self.embedding_model = embeddings
        self.chunk_size = 1000
        self.chunk_overlap = 200
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            is_separator_regex=False,
        )
        
        # Supported file extensions
        self.supported_extensions = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt'}
    
    def _get_local_files(self, directory_path: str, recursive: bool = True) -> List[Dict[str, Any]]:
        """Get list of supported documents from local directory"""
        documents = []
        path = Path(directory_path)
        
        # Get files recursively or from single directory
        if recursive:
            file_pattern = '**/*'
        else:
            file_pattern = '*'
        
        for file_path in path.glob(file_pattern):
            if file_path.is_file():
                file_extension = file_path.suffix.lower()
                
                if file_extension in self.supported_extensions:
                    # Get file metadata
                    stat_info = file_path.stat()
                    
                    documents.append({
                        'id': str(hash(str(file_path.absolute()))),  # Generate unique ID from path
                        'name': file_path.name,
                        'extension': file_extension,
                        'path': str(file_path.absolute()),
                        'relative_path': str(file_path.relative_to(path)),
                        'size': stat_info.st_size,
                        'modified': stat_info.st_mtime,
                        'parent_directory': file_path.parent.name
                    })
                    print(f"Opening File : {file_path.name}")

        return documents
    
    def _read_file(self, file_path: str) -> bytes:
        """Read file content from local filesystem"""
        try:
            with open(file_path, 'rb') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading file {file_path}: {e}")
            return b''
    
    def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF using PyMuPDF (fitz)"""
        text = ""
        try:
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text += page.get_text()
            
            pdf_document.close()
            print(f"Extracted text from PDF")
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            text = ""
        
        return text
    
    def _extract_word_text(self, file_content: bytes) -> str:
        """Extract text from Word documents"""
        text = ""
        try:
            doc = Document(io.BytesIO(file_content))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                text += "\n"
        except Exception as e:
            print(f"Error extracting Word text: {e}")
            text = ""
        
        return text
    
    def _extract_powerpoint_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        text = ""
        try:
            prs = Presentation(io.BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting PowerPoint text: {e}")
            text = ""
        
        return text
    
    def _extract_text_file(self, file_content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            return file_content.decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"Error extracting text: {e}")
            return ""
    
    
    async def ingest_documents(self, directory_path: str, recursive: bool = True) -> Dict[str, Any]:
        documents = self._get_local_files(directory_path, recursive)

        if not documents:
            return {
                "total_documents": 0,
                "indexed_count": 0,
                "errors": ["No supported documents found"]
            }

        # Run document processing concurrently
        tasks = [self._process_document(doc) for doc in documents]
        results = await asyncio.gather(*tasks, return_exceptions=True)

        errors = [str(e) for e in results if isinstance(e, Exception)]
        valid_results = [r for r in results if not isinstance(r, Exception)]

        indexed_count = await self._index_documents(valid_results)

        return {
            "total_documents": len(documents),
            "indexed_count": indexed_count,
            "errors": errors
        }

    
    async def _process_document(self, doc_info: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Process individual document"""
        try:
            file_content = self._read_file(doc_info['path'])
            
            if not file_content:
                return []
            
            # Extract text based on file type
            if doc_info['extension'] == '.pdf':
                text = self._extract_pdf_text(file_content)
            elif doc_info['extension'] in ['.docx', '.doc']:
                text = self._extract_word_text(file_content)
            elif doc_info['extension'] in ['.pptx', '.ppt']:
                text = self._extract_powerpoint_text(file_content)
            elif doc_info['extension'] == '.txt':
                text = self._extract_text_file(file_content)
            else:
                text = file_content.decode('utf-8', errors='ignore')
            
            # Skip empty documents
            if not text.strip():
                print(f"Warning: No text extracted from {doc_info['name']}")
                return []
            
            # Chunk text
            chunks = self._chunk_text(text)
            
            # Generate embeddings and summaries for each chunk
            processed_chunks = []
            for i, chunk in enumerate(chunks):
                
                processed_chunks.append({
                    'doc_id': f"{doc_info['id']}_{i}",
                    'source': 'LocalFile',
                    'file_type': doc_info['extension'],
                    'page_number': i,
                    'title': doc_info['name'],
                    'content': chunk,
                    'metadata': {
                        'file_path': doc_info['path'],
                        'relative_path': doc_info['relative_path'],
                        'parent_directory': doc_info['parent_directory'],
                        'file_size': doc_info['size'],
                        'modified': doc_info['modified'],
                        'chunk_index': i,
                        'total_chunks': len(chunks)
                    }
                })
            
            return processed_chunks
        except Exception as e:
            print(f"Error processing document {doc_info.get('name', 'unknown')}: {e}")
            return []
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks"""
        return self.text_splitter.split_text(text)
    
    async def _index_documents(self, document_chunks: List[List[Dict]]) -> int:
        """Index processed documents to Azure AI Search with proper bulk upload"""
        # Flatten the list of chunks
        all_chunks = []
        for doc_chunks in document_chunks:
            if isinstance(doc_chunks, list):
                all_chunks.extend(doc_chunks)
        
        if not all_chunks:
            return 0
        
        # Prepare texts and metadata for bulk upload
        texts = []
        metadatas = []
        keys = []
        
        for chunk in all_chunks:
            texts.append(chunk['content'])
            
            # Prepare metadata dictionary
            metadata = {
                'source': chunk['source'],
                'file_type': chunk['file_type'],
                'page_number': chunk['page_number'],
                'title': chunk['title'],
                'file_path': chunk['metadata'].get('file_path', ''),
                'relative_path': chunk['metadata'].get('relative_path', ''),
                'parent_directory': chunk['metadata'].get('parent_directory', ''),
                'file_size': chunk['metadata'].get('file_size', 0),
                'modified': chunk['metadata'].get('modified', ''),
                'chunk_index': chunk['metadata'].get('chunk_index', 0),
                'total_chunks': chunk['metadata'].get('total_chunks', 0)
            }
            metadatas.append(metadata)
            keys.append(chunk['doc_id'])
        
        try:
            # Use add_texts method with proper parameters
            ids = self.azure_ai_search_client.add_texts(
                texts=texts,
                metadatas=metadatas,
                keys=keys
            )
            
            return len(ids)
        except Exception as e:
            print(f"Error indexing documents to Azure AI Search: {e}")
            return 0


# Example usage
async def main():
    # Initialize the ingestion class
    ingestion = LocalFileIngestion(index_name="local-documents")
    
    # Specify the local directory path containing documents
    local_directory = "./documents"  # Change to your local directory path
    
    # Ingest documents (recursive=True will scan subdirectories)
    result = await ingestion.ingest_documents(
        directory_path=local_directory,
        recursive=True  # Set to False to only scan the top-level directory
    )
    
    print(f"Total documents: {result['total_documents']}")
    print(f"Indexed documents: {result['indexed_count']}")
    if result['errors']:
        print(f"Errors: {result['errors']}")


if __name__ == "__main__":
    asyncio.run(main())


# Example usage
async def main():
    ingestion = SharePointIngestion(index_name="sharepoint-documents")
    
    # Specify the SharePoint document library path
    # Format: "/sites/YourSite/Shared Documents" or "/Shared Documents"
    document_library_path = "/Shared Documents"
    
    result = await ingestion.ingest_documents(document_library_path)
    
    print(f"Total documents: {result['total_documents']}")
    print(f"Indexed documents: {result['indexed_count']}")
    if result['errors']:
        print(f"Errors: {result['errors']}")


if __name__ == "__main__":
    asyncio.run(main())
